"""
Generate one test file per format supported by the MarkItDown converter.
Run once:  python generate_test_files.py
Output goes into test-files/
"""
import os, json, pathlib

OUT = pathlib.Path("test-files")
OUT.mkdir(exist_ok=True)

# ── HTML ──────────────────────────────────────────────────────────────────────
html = OUT / "sample.html"
html.write_text("""\
<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Sample HTML Page</title></head>
<body>
  <h1>MarkItDown HTML Test</h1>
  <p>This is a <strong>bold</strong> statement and this is <em>italic</em>.</p>
  <h2>Features</h2>
  <ul>
    <li>Convert HTML to Markdown</li>
    <li>Preserves headings and lists</li>
    <li>Handles <a href="https://example.com">hyperlinks</a></li>
  </ul>
  <h2>Data Table</h2>
  <table>
    <tr><th>Name</th><th>Role</th><th>Score</th></tr>
    <tr><td>Alice</td><td>Engineer</td><td>95</td></tr>
    <tr><td>Bob</td><td>Designer</td><td>88</td></tr>
    <tr><td>Carol</td><td>PM</td><td>91</td></tr>
  </table>
  <blockquote>The best way to predict the future is to invent it. — Alan Kay</blockquote>
</body>
</html>
""", encoding="utf-8")
print("✓ sample.html")

# ── CSV ───────────────────────────────────────────────────────────────────────
csv = OUT / "sample.csv"
csv.write_text("""\
Product,Category,Price,Stock,Rating
Wireless Mouse,Electronics,29.99,150,4.5
Standing Desk,Furniture,349.00,42,4.8
Coffee Maker,Appliances,89.99,78,4.3
Notebook Set,Stationery,12.49,300,4.6
Ergonomic Chair,Furniture,499.00,25,4.9
USB-C Hub,Electronics,49.99,200,4.4
""", encoding="utf-8")
print("✓ sample.csv")

# ── JSON ──────────────────────────────────────────────────────────────────────
data = {
    "project": "MarkItDown Test",
    "version": "1.0.0",
    "description": "Sample JSON file for converter testing",
    "team": [
        {"name": "Alice", "role": "Backend", "languages": ["Python", "Go"]},
        {"name": "Bob",   "role": "Frontend", "languages": ["TypeScript", "CSS"]},
        {"name": "Carol", "role": "DevOps",   "languages": ["Bash", "YAML"]},
    ],
    "config": {
        "maxFileSize": "50MB",
        "supportedFormats": ["pdf", "docx", "xlsx", "pptx", "html", "csv", "json", "txt"],
        "deployed": True,
    },
}
(OUT / "sample.json").write_text(json.dumps(data, indent=2), encoding="utf-8")
print("✓ sample.json")

# ── PNG image ─────────────────────────────────────────────────────────────────
from PIL import Image, ImageDraw, ImageFont

img = Image.new("RGB", (600, 300), color=(30, 30, 46))
draw = ImageDraw.Draw(img)
# Background gradient stripes
for i in range(0, 600, 40):
    draw.rectangle([i, 0, i + 20, 300], fill=(45, 45, 65))
# Text
draw.rectangle([50, 50, 550, 250], fill=(50, 50, 70), outline=(100, 200, 100), width=2)
draw.text((300, 100), "MarkItDown", fill=(100, 220, 100), anchor="mm")
draw.text((300, 145), "Image Conversion Test", fill=(180, 180, 200), anchor="mm")
draw.text((300, 190), "PNG · 600 × 300 px", fill=(120, 120, 150), anchor="mm")
img.save(OUT / "sample.png")
print("✓ sample.png")

# ── XLSX (Excel) ──────────────────────────────────────────────────────────────
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment

wb = openpyxl.Workbook()

# Sheet 1 – Sales data
ws1 = wb.active
ws1.title = "Q1 Sales"
headers = ["Month", "Region", "Product", "Units", "Revenue ($)"]
ws1.append(headers)
for cell in ws1[1]:
    cell.font = Font(bold=True, color="FFFFFF")
    cell.fill = PatternFill("solid", fgColor="2E7D32")
    cell.alignment = Alignment(horizontal="center")
rows = [
    ("January",  "North", "Widget A", 320, 9600),
    ("January",  "South", "Widget B", 210, 8400),
    ("February", "North", "Widget A", 410, 12300),
    ("February", "East",  "Widget C", 175, 8750),
    ("March",    "West",  "Widget B", 290, 11600),
    ("March",    "North", "Widget C", 330, 16500),
]
for row in rows:
    ws1.append(row)
ws1.column_dimensions["A"].width = 12
ws1.column_dimensions["B"].width = 10
ws1.column_dimensions["C"].width = 14
ws1.column_dimensions["D"].width = 8
ws1.column_dimensions["E"].width = 14

# Sheet 2 – Summary
ws2 = wb.create_sheet("Summary")
ws2.append(["Metric", "Value"])
ws2.append(["Total Units", "=SUM('Q1 Sales'!D2:D7)"])
ws2.append(["Total Revenue", "=SUM('Q1 Sales'!E2:E7)"])
ws2.append(["Avg Revenue/Unit", "=B3/B2"])

wb.save(OUT / "sample.xlsx")
print("✓ sample.xlsx")

# ── PPTX (PowerPoint) ─────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
W, H = prs.slide_width, prs.slide_height

# Slide 1 – Title
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "MarkItDown Converter"
slide1.placeholders[1].text = "PowerPoint → Markdown Test File"

# Slide 2 – Bullet list
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Supported Formats"
tf = slide2.placeholders[1].text_frame
tf.text = "Document Formats"
for item in ["PDF files", "Word documents (.docx)", "Excel spreadsheets (.xlsx)"]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# Slide 3 – Text content
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "About This Test"
tf3 = slide3.placeholders[1].text_frame
tf3.text = "This PPTX was generated to verify that the MarkItDown converter correctly extracts text from PowerPoint files and renders it as Markdown."

prs.save(OUT / "sample.pptx")
print("✓ sample.pptx")

# ── DOCX (Word) ──────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor

doc = Document()
doc.add_heading("MarkItDown Word Document Test", 0)
doc.add_paragraph(
    "This document was generated to verify that the MarkItDown converter "
    "correctly extracts text from Word files and renders it as Markdown."
)
doc.add_heading("Section 1: Formatting", level=1)
p = doc.add_paragraph()
p.add_run("Bold text").bold = True
p.add_run(" and ")
p.add_run("italic text").italic = True
p.add_run(" and ")
p.add_run("underlined text").underline = True
p.add_run(" all in one paragraph.")

doc.add_heading("Section 2: Lists", level=1)
doc.add_paragraph("Unordered items:", style="List Bullet")
for item in ["First bullet point", "Second bullet point", "Third bullet point"]:
    doc.add_paragraph(item, style="List Bullet")

doc.add_heading("Section 3: Table", level=1)
table = doc.add_table(rows=4, cols=3)
table.style = "Table Grid"
hdr = table.rows[0].cells
hdr[0].text, hdr[1].text, hdr[2].text = "Name", "Department", "Years"
data_rows = [("Alice", "Engineering", "5"), ("Bob", "Design", "3"), ("Carol", "Product", "7")]
for i, (name, dept, yrs) in enumerate(data_rows, 1):
    row = table.rows[i].cells
    row[0].text, row[1].text, row[2].text = name, dept, yrs

doc.add_heading("Section 4: Quote", level=1)
doc.add_paragraph(
    '"The best tool is the one that converts your files correctly." — MarkItDown',
    style="Intense Quote",
)
doc.save(OUT / "sample.docx")
print("✓ sample.docx")

# ── PDF ───────────────────────────────────────────────────────────────────────
from fpdf import FPDF

pdf = FPDF()
pdf.set_auto_page_break(auto=True, margin=15)

# Page 1
pdf.add_page()
pdf.set_font("Helvetica", "B", 24)
pdf.cell(0, 14, "MarkItDown PDF Test", ln=True, align="C")
pdf.ln(4)
pdf.set_font("Helvetica", size=12)
pdf.multi_cell(0, 8,
    "This PDF was generated to verify that the MarkItDown converter correctly "
    "extracts text from PDF files and renders it as Markdown."
)
pdf.ln(6)

pdf.set_font("Helvetica", "B", 16)
pdf.cell(0, 10, "Section 1: Introduction", ln=True)
pdf.set_font("Helvetica", size=12)
pdf.multi_cell(0, 8,
    "MarkItDown is an open-source library by Microsoft that converts a wide variety "
    "of file formats into Markdown. It supports PDFs, Word documents, Excel sheets, "
    "PowerPoint presentations, images, HTML pages, CSV files, and JSON data."
)
pdf.ln(4)

pdf.set_font("Helvetica", "B", 16)
pdf.cell(0, 10, "Section 2: Key Features", ln=True)
pdf.set_font("Helvetica", size=12)
for feature in [
    "- Fast, lightweight conversion",
    "- Preserves document structure (headings, lists, tables)",
    "- Handles multi-page documents",
    "- Zero external API calls required",
    "- Open source (MIT license)",
]:
    pdf.cell(0, 8, feature, ln=True)
pdf.ln(4)

# Page 2
pdf.add_page()
pdf.set_font("Helvetica", "B", 16)
pdf.cell(0, 10, "Section 3: Sample Table", ln=True)
pdf.set_font("Helvetica", "B", 11)
col_w = [50, 50, 45, 45]
for header, w in zip(["Name", "Role", "Language", "Score"], col_w):
    pdf.cell(w, 8, header, border=1)
pdf.ln()
pdf.set_font("Helvetica", size=11)
for row in [
    ("Alice", "Backend",  "Python",     "95"),
    ("Bob",   "Frontend", "TypeScript", "88"),
    ("Carol", "DevOps",   "Bash",       "91"),
    ("Dave",  "Data",     "SQL",        "84"),
]:
    for val, w in zip(row, col_w):
        pdf.cell(w, 8, val, border=1)
    pdf.ln()

pdf.ln(6)
pdf.set_font("Helvetica", "I", 11)
pdf.multi_cell(0, 8,
    "Quote: \"Any sufficiently advanced technology is indistinguishable from magic.\" "
    "- Arthur C. Clarke"
)

pdf.output(str(OUT / "sample.pdf"))
print("✓ sample.pdf")

print(f"\nAll 8 test files written to ./{OUT}/")
for f in sorted(OUT.iterdir()):
    print(f"  {f.name:30s}  {f.stat().st_size:>10,} bytes")
